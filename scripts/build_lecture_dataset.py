"""
Build the LectureSlideOCR-500-DE benchmark dataset.
Phase 4 (Week 10-11): Create the novel dataset for publication.

LectureSlideOCR-500-DE:
    - 500 annotated lecture slide images from German-speaking professors
    - Handwritten German text + mathematical expressions
    - Bounding box annotations (text/math regions)
    - Ground-truth transcriptions (verified)
    - Professor-specific splits for few-shot evaluation

Dataset structure:
    data/processed/lecture_slides/
    ├── images/
    │   ├── train/ (350 slides)
    │   ├── val/   (75 slides)
    │   └── test/  (75 slides)
    ├── annotations/
    │   ├── train.json   (COCO format)
    │   ├── val.json
    │   └── test.json
    ├── transcriptions/
    │   ├── train.json
    │   ├── val.json
    │   └── test.json
    ├── professor_splits/
    │   ├── professor_{id}/
    │   │   ├── support.json  (50 samples)
    │   │   └── query.json    (100+ samples)
    └── dataset_info.json

Usage:
    # From raw slide images + annotation CSV/JSON:
    python scripts/build_lecture_dataset.py \
        --source data/Dr_Judith_Jakob_Slides/ \
        --output data/processed/lecture_slides/ \
        --annotator human \
        --professor-id dr_jakob

    # Auto-annotate with current pipeline (for bootstrapping):
    python scripts/build_lecture_dataset.py \
        --source data/Dr_Judith_Jakob_Slides/ \
        --output data/processed/lecture_slides/ \
        --annotator pipeline \
        --detector runs/detect/baseline_v1/weights/best.pt \
        --professor-id dr_jakob

Note:
    DO NOT run this on Dr_Judith_Jakob_Slides/ without user confirmation.
    The professor slides require explicit permission before processing.
    This script provides the infrastructure ready for when data is available.
"""

from __future__ import annotations

import argparse
import json
import shutil
import sys
import time
from pathlib import Path
from typing import List, Dict, Optional, Tuple

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

import numpy as np
from loguru import logger


# ---------------------------------------------------------------------------
# Dataset schema
# ---------------------------------------------------------------------------

def create_dataset_info(
    professor_id: str,
    n_train: int,
    n_val: int,
    n_test: int,
    source_desc: str,
) -> dict:
    """Create dataset_info.json for the benchmark."""
    return {
        'name': 'LectureSlideOCR-500-DE',
        'version': '1.0',
        'description': (
            'German lecture slide OCR benchmark with handwritten text and math. '
            'Contains slides from German-speaking professors with bounding box '
            'annotations and verified transcriptions.'
        ),
        'classes': ['text', 'math'],
        'class_ids': {'text': 0, 'math': 1},
        'splits': {
            'train': n_train,
            'val': n_val,
            'test': n_test,
        },
        'professors': [professor_id],
        'source': source_desc,
        'created': time.strftime('%Y-%m-%d'),
        'citation': {
            'title': 'LectureSlideOCR-500-DE: A Benchmark for German Handwriting OCR '
                     'in Lecture Slides with Meta-Learning Adaptation',
            'venue': 'ICDAR 2026 / CVPR 2026',
        }
    }


def coco_annotation_template(image_id: int, image_path: str, height: int, width: int) -> dict:
    """COCO-format image entry."""
    return {
        'id': image_id,
        'file_name': Path(image_path).name,
        'height': height,
        'width': width,
        'professor_id': None,
        'slide_source': str(image_path),
    }


def coco_bbox_annotation(
    ann_id: int,
    image_id: int,
    category_id: int,
    bbox: Tuple[float, float, float, float],  # x, y, w, h
    text: str = '',
) -> dict:
    """COCO-format annotation entry."""
    x, y, w, h = bbox
    return {
        'id': ann_id,
        'image_id': image_id,
        'category_id': category_id,
        'bbox': [x, y, w, h],
        'area': w * h,
        'iscrowd': 0,
        'text': text,
        'verified': False,
    }


# ---------------------------------------------------------------------------
# Image extraction
# ---------------------------------------------------------------------------

def extract_slide_images(
    source_dir: Path,
    output_dir: Path,
    extensions: Tuple[str, ...] = ('.png', '.jpg', '.jpeg', '.bmp', '.tiff'),
) -> List[Path]:
    """
    Extract slide images from a directory.

    Supports image files and PDF extraction (requires pdf2image).
    """
    output_dir.mkdir(parents=True, exist_ok=True)
    images = []

    # Direct image files
    for ext in extensions:
        images.extend(sorted(source_dir.rglob(f'*{ext}')))
        images.extend(sorted(source_dir.rglob(f'*{ext.upper()}')))

    # PDF files
    pdf_files = list(source_dir.rglob('*.pdf'))
    if pdf_files:
        logger.info(f"Found {len(pdf_files)} PDF files, extracting pages...")
        for pdf in pdf_files:
            try:
                from pdf2image import convert_from_path
                pages = convert_from_path(str(pdf), dpi=150)
                for i, page in enumerate(pages):
                    out_path = output_dir / f"{pdf.stem}_page{i+1:03d}.png"
                    page.save(str(out_path))
                    images.append(out_path)
                logger.info(f"  {pdf.name}: {len(pages)} pages")
            except ImportError:
                logger.warning("pdf2image not installed. Install: pip install pdf2image poppler")
            except Exception as e:
                logger.warning(f"Failed to extract {pdf}: {e}")

    # PPTX files
    pptx_files = list(source_dir.rglob('*.pptx'))
    if pptx_files:
        logger.info(f"Found {len(pptx_files)} PPTX files, extracting slides...")
        for pptx in pptx_files:
            extracted = _extract_pptx_images(pptx, output_dir)
            images.extend(extracted)

    logger.info(f"Found {len(images)} slide images in {source_dir}")
    return sorted(set(images))


def _extract_pptx_images(pptx_path: Path, output_dir: Path) -> List[Path]:
    """Extract PPTX slides as images (requires python-pptx + Pillow)."""
    extracted = []
    try:
        from pptx import Presentation
        from pptx.util import Pt
        from PIL import Image as PILImage
        import io

        prs = Presentation(str(pptx_path))
        for i, slide in enumerate(prs.slides):
            # Extract embedded images
            for shape in slide.shapes:
                if shape.shape_type == 13:  # Picture
                    image = shape.image
                    image_bytes = image.blob
                    pil = PILImage.open(io.BytesIO(image_bytes))
                    out_path = output_dir / f"{pptx_path.stem}_slide{i+1:03d}_img.png"
                    pil.save(str(out_path))
                    extracted.append(out_path)

        logger.info(f"Extracted {len(extracted)} images from {pptx_path.name}")
    except ImportError:
        logger.warning("python-pptx not installed. Install: pip install python-pptx")
    except Exception as e:
        logger.warning(f"PPTX extraction failed for {pptx_path}: {e}")

    return extracted


# ---------------------------------------------------------------------------
# Auto-annotation with pipeline
# ---------------------------------------------------------------------------

def auto_annotate_with_pipeline(
    image_paths: List[Path],
    detector_path: str,
    output_dir: Path,
    conf: float = 0.35,
    device: str = 'cuda',
) -> Tuple[dict, dict]:
    """
    Bootstrap annotations using the trained detector.

    Returns COCO-format annotation dict. All detections marked as unverified.
    Human review required before use in evaluation.

    Args:
        image_paths: Slide images to annotate.
        detector_path: YOLOv8 checkpoint path.
        output_dir: Where to save annotations.
        conf: Detection confidence threshold.
        device: Compute device.

    Returns:
        Tuple of (coco_annotations, transcriptions)
    """
    from baseline.baseline_pipeline import YOLOv8Detector, TrOCRRecognizer, BaselinePipeline
    from utils.image_utils import load_image

    logger.info(f"Auto-annotating {len(image_paths)} slides with pipeline...")
    logger.warning("NOTE: Auto-annotations require human verification before evaluation use!")

    detector = YOLOv8Detector(model_path=detector_path, conf_threshold=conf, device=device)
    text_ocr = TrOCRRecognizer(device=device)

    coco = {
        'images': [],
        'annotations': [],
        'categories': [
            {'id': 0, 'name': 'text', 'supercategory': 'text'},
            {'id': 1, 'name': 'math', 'supercategory': 'math'},
        ],
    }
    transcriptions = {'samples': []}

    ann_id = 0
    for img_id, img_path in enumerate(image_paths):
        try:
            img = load_image(img_path, mode='rgb')
            h, w = img.shape[:2]

            coco['images'].append(
                coco_annotation_template(img_id, str(img_path), h, w)
            )

            # Detect regions
            detections = detector.detect(img)
            img_anns = []

            for det in detections:
                x1, y1, x2, y2 = det['bbox']
                bw, bh = x2 - x1, y2 - y1
                cat_id = 0 if det['class'] == 'text' else 1

                # Crop and OCR if text
                text = ''
                if det['class'] == 'text':
                    try:
                        from PIL import Image as PILImage
                        crop = img[int(y1):int(y2), int(x1):int(x2)]
                        if crop.size > 0:
                            pil_crop = PILImage.fromarray(crop)
                            text = text_ocr.recognize(pil_crop) or ''
                    except Exception:
                        pass

                ann = coco_bbox_annotation(ann_id, img_id, cat_id, [x1, y1, bw, bh], text)
                coco['annotations'].append(ann)
                img_anns.append(ann)
                ann_id += 1

                if text:
                    transcriptions['samples'].append({
                        'image': str(img_path),
                        'bbox': [x1, y1, bw, bh],
                        'text': text,
                        'class': det['class'],
                        'confidence': det.get('conf', 0.0),
                        'verified': False,
                    })

            if img_id % 50 == 0:
                logger.info(f"Progress: {img_id+1}/{len(image_paths)} slides, "
                            f"{ann_id} annotations")

        except Exception as e:
            logger.warning(f"Failed to annotate {img_path}: {e}")

    logger.info(f"Auto-annotation complete: {len(coco['images'])} images, "
                f"{len(coco['annotations'])} annotations")
    return coco, transcriptions


# ---------------------------------------------------------------------------
# Dataset split and professor partition
# ---------------------------------------------------------------------------

def create_splits(
    image_paths: List[Path],
    train_ratio: float = 0.7,
    val_ratio: float = 0.15,
    test_ratio: float = 0.15,
    seed: int = 42,
) -> Tuple[List[Path], List[Path], List[Path]]:
    """Split images into train/val/test."""
    rng = np.random.RandomState(seed)
    indices = rng.permutation(len(image_paths))

    n_train = int(len(image_paths) * train_ratio)
    n_val = int(len(image_paths) * val_ratio)

    train_idx = indices[:n_train]
    val_idx = indices[n_train:n_train + n_val]
    test_idx = indices[n_train + n_val:]

    return (
        [image_paths[i] for i in train_idx],
        [image_paths[i] for i in val_idx],
        [image_paths[i] for i in test_idx],
    )


def create_professor_split(
    all_annotations: List[Dict],
    professor_id: str,
    n_support: int = 50,
    output_dir: Path = Path('data/processed/lecture_slides/professor_splits'),
) -> None:
    """
    Create professor-specific support/query split for few-shot evaluation.

    Support set: n_support verified text transcriptions
    Query set: remaining samples
    """
    prof_dir = output_dir / f'professor_{professor_id}'
    prof_dir.mkdir(parents=True, exist_ok=True)

    text_samples = [a for a in all_annotations if a.get('class') == 'text' and a.get('text')]
    rng = np.random.RandomState(42)
    rng.shuffle(text_samples)

    support = {'professor_id': professor_id, 'samples': text_samples[:n_support]}
    query = {'professor_id': professor_id, 'samples': text_samples[n_support:]}

    with open(prof_dir / 'support.json', 'w', encoding='utf-8') as f:
        json.dump(support, f, ensure_ascii=False, indent=2)
    with open(prof_dir / 'query.json', 'w', encoding='utf-8') as f:
        json.dump(query, f, ensure_ascii=False, indent=2)

    logger.info(f"Professor split for '{professor_id}': "
                f"{len(support['samples'])} support, {len(query['samples'])} query")


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def build_lecture_dataset(
    source_dir: Path,
    output_dir: Path,
    professor_id: str,
    annotator: str = 'human',
    detector_path: Optional[str] = None,
    n_support: int = 50,
    device: str = 'cuda',
) -> None:
    """
    Build the LectureSlideOCR-500-DE benchmark from raw slide images.

    Args:
        source_dir: Directory containing raw slide images / PDFs / PPTXs.
        output_dir: Output directory for processed dataset.
        professor_id: Identifier for the professor (e.g., 'dr_jakob').
        annotator: 'human' (use existing annotation files) or 'pipeline' (auto-annotate).
        detector_path: YOLOv8 checkpoint for auto-annotation.
        n_support: Number of professor-specific support samples.
        device: Compute device.
    """
    logger.info("=" * 60)
    logger.info("Building LectureSlideOCR-500-DE Dataset")
    logger.info(f"  Source:     {source_dir}")
    logger.info(f"  Output:     {output_dir}")
    logger.info(f"  Professor:  {professor_id}")
    logger.info(f"  Annotator:  {annotator}")
    logger.info("=" * 60)

    output_dir = Path(output_dir)
    images_dir = output_dir / 'images'
    anns_dir = output_dir / 'annotations'
    trans_dir = output_dir / 'transcriptions'

    for d in [images_dir, anns_dir, trans_dir]:
        d.mkdir(parents=True, exist_ok=True)

    # Extract images
    logger.info("Step 1: Extracting slide images...")
    all_images = extract_slide_images(source_dir, images_dir / 'raw')

    if not all_images:
        logger.error(f"No images found in {source_dir}")
        return

    # Split
    logger.info("Step 2: Creating train/val/test splits...")
    train_imgs, val_imgs, test_imgs = create_splits(all_images)
    logger.info(f"  Train: {len(train_imgs)}, Val: {len(val_imgs)}, Test: {len(test_imgs)}")

    # Copy to split directories
    for split_name, split_imgs in [('train', train_imgs), ('val', val_imgs), ('test', test_imgs)]:
        split_dir = images_dir / split_name
        split_dir.mkdir(exist_ok=True)
        for img_path in split_imgs:
            shutil.copy2(str(img_path), str(split_dir / img_path.name))

    # Annotate
    if annotator == 'pipeline':
        if not detector_path:
            logger.error("--detector-path required for pipeline annotation")
            return
        logger.info("Step 3: Auto-annotating with pipeline (human verification needed)...")
        coco_anns, transcriptions = auto_annotate_with_pipeline(
            all_images, detector_path, anns_dir, device=device
        )
        all_transcription_samples = transcriptions['samples']
    else:
        logger.info("Step 3: Using human annotations (expected in source_dir)...")
        # Look for existing annotation files
        coco_file = source_dir / 'annotations.json'
        trans_file = source_dir / 'transcriptions.json'

        if coco_file.exists():
            with open(coco_file) as f:
                coco_anns = json.load(f)
            logger.info(f"Loaded {len(coco_anns.get('annotations', []))} human annotations")
        else:
            logger.warning(f"No annotations.json found in {source_dir}")
            logger.warning("Creating empty annotation structure for manual filling")
            coco_anns = {
                'images': [],
                'annotations': [],
                'categories': [
                    {'id': 0, 'name': 'text'},
                    {'id': 1, 'name': 'math'},
                ],
            }

        if trans_file.exists():
            with open(trans_file) as f:
                transcriptions = json.load(f)
            all_transcription_samples = transcriptions.get('samples', [])
        else:
            all_transcription_samples = []

    # Save annotations per split
    logger.info("Step 4: Saving split annotations...")
    for split_name, split_imgs in [('train', train_imgs), ('val', val_imgs), ('test', test_imgs)]:
        split_names = {img.name for img in split_imgs}
        split_img_ids = {
            img['id'] for img in coco_anns.get('images', [])
            if img.get('file_name') in split_names
        }

        split_coco = {
            'images': [
                img for img in coco_anns.get('images', []) if img['id'] in split_img_ids
            ],
            'annotations': [
                ann for ann in coco_anns.get('annotations', []) if ann['image_id'] in split_img_ids
            ],
            'categories': coco_anns.get('categories', []),
        }

        with open(anns_dir / f'{split_name}.json', 'w', encoding='utf-8') as f:
            json.dump(split_coco, f, ensure_ascii=False, indent=2)
        logger.info(f"  {split_name}: {len(split_coco['images'])} images, "
                    f"{len(split_coco['annotations'])} annotations")

        # Save transcriptions for this split
        split_trans = {
            'samples': [
                s for s in all_transcription_samples
                if Path(s.get('image', '')).name in split_names
            ]
        }
        with open(trans_dir / f'{split_name}.json', 'w', encoding='utf-8') as f:
            json.dump(split_trans, f, ensure_ascii=False, indent=2)

    # Create professor-specific split
    logger.info("Step 5: Creating professor-specific few-shot split...")
    create_professor_split(
        all_annotations=all_transcription_samples,
        professor_id=professor_id,
        n_support=n_support,
        output_dir=output_dir / 'professor_splits',
    )

    # Dataset info
    dataset_info = create_dataset_info(
        professor_id=professor_id,
        n_train=len(train_imgs),
        n_val=len(val_imgs),
        n_test=len(test_imgs),
        source_desc=str(source_dir),
    )
    with open(output_dir / 'dataset_info.json', 'w') as f:
        json.dump(dataset_info, f, indent=2)

    logger.info("=" * 60)
    logger.info("Dataset built successfully!")
    logger.info(f"  Total images:  {len(all_images)}")
    logger.info(f"  Annotations:   {len(coco_anns.get('annotations', []))}")
    logger.info(f"  Transcriptions:{len(all_transcription_samples)}")
    logger.info(f"  Output:        {output_dir}")
    logger.info("=" * 60)


def parse_args():
    parser = argparse.ArgumentParser(
        description='Build LectureSlideOCR-500-DE benchmark',
        epilog=(
            'IMPORTANT: Do not run on Dr_Judith_Jakob_Slides/ without user confirmation. '
            'Professor slides require explicit permission before processing.'
        )
    )
    parser.add_argument('--source', type=Path, required=True,
                        help='Source directory with slide images/PDFs/PPTXs')
    parser.add_argument('--output', type=Path,
                        default=Path('data/processed/lecture_slides'),
                        help='Output directory for processed dataset')
    parser.add_argument('--professor-id', type=str, required=True,
                        help='Professor identifier (e.g., dr_jakob)')
    parser.add_argument('--annotator', type=str, default='human',
                        choices=['human', 'pipeline'],
                        help='Annotation source: human (existing files) or pipeline (auto)')
    parser.add_argument('--detector', type=str, default=None,
                        help='YOLOv8 checkpoint for auto-annotation')
    parser.add_argument('--n-support', type=int, default=50,
                        help='Professor support samples for few-shot evaluation')
    parser.add_argument('--device', type=str, default='cuda')
    return parser.parse_args()


if __name__ == '__main__':
    args = parse_args()

    # Safety check
    source_str = str(args.source).replace('\\', '/')
    if 'Jakob' in source_str or 'jakob' in source_str or 'Dr_Judith' in source_str:
        logger.warning("=" * 60)
        logger.warning("WARNING: You are about to process professor slides!")
        logger.warning(f"Source: {args.source}")
        logger.warning("This requires explicit permission from the professor.")
        logger.warning("Are you sure you want to proceed? (yes/no)")
        logger.warning("=" * 60)
        confirm = input("Confirm (yes/no): ").strip().lower()
        if confirm != 'yes':
            logger.info("Aborted.")
            sys.exit(0)

    build_lecture_dataset(
        source_dir=args.source,
        output_dir=args.output,
        professor_id=args.professor_id,
        annotator=args.annotator,
        detector_path=args.detector,
        n_support=args.n_support,
        device=args.device,
    )
